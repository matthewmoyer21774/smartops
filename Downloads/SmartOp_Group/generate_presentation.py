"""
Generate a PowerPoint presentation for the inventory management methodology.
Run: python generate_presentation.py
Output: Presentation.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Color palette ---
DARK_BLUE = RGBColor(0, 70, 127)
MED_BLUE = RGBColor(0, 120, 190)
LIGHT_BLUE = RGBColor(0, 170, 220)
ACCENT_ORANGE = RGBColor(230, 130, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(60, 60, 60)
MED_GRAY = RGBColor(120, 120, 120)
LIGHT_BG = RGBColor(240, 245, 250)
RED = RGBColor(200, 50, 50)
GREEN = RGBColor(40, 150, 70)


def add_bg_rect(slide, color=DARK_BLUE):
    """Add a full-width top bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_footer_bar(slide):
    """Add a thin bottom bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.1), prs.slide_width, Inches(0.4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.text = "SmartOp Final Assignment  |  KU Leuven"
    for p in tf.paragraphs:
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER


def add_slide_title(slide, title_text, subtitle_text=None):
    """Add title text in the top bar area."""
    add_bg_rect(slide)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(180, 210, 240)

    add_footer_bar(slide)


def add_body_text(slide, left, top, width, height, bullets, font_size=18, bold_first_word=False):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        p.space_before = Pt(4)

        if isinstance(bullet, tuple):
            # (bold_part, normal_part)
            run1 = p.add_run()
            run1.text = bullet[0]
            run1.font.bold = True
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = DARK_GRAY
            run2 = p.add_run()
            run2.text = bullet[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = DARK_GRAY
        else:
            p.text = bullet
            p.font.size = Pt(font_size)
            p.font.color.rgb = DARK_GRAY

    return txBox


def add_table(slide, left, top, width, height, data, col_widths=None):
    """Add a table to the slide."""
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.CENTER
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_GRAY

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG

    return table_shape


def add_accent_box(slide, left, top, width, height, text, bg_color=LIGHT_BG, font_color=DARK_BLUE, font_size=20):
    """Add a colored accent box with text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = MED_BLUE
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = True
    return shape


# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Full blue background
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Data-Driven Inventory Management"
p.font.size = Pt(44)
p.font.color.rgb = WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Perishable SKU Ordering Policy"
p2.font.size = Pt(28)
p2.font.color.rgb = LIGHT_BLUE
p2.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10), Inches(2))
tf2 = txBox2.text_frame
p3 = tf2.paragraphs[0]
p3.text = "SmartOp - Final Assignment"
p3.font.size = Pt(22)
p3.font.color.rgb = RGBColor(180, 210, 240)
p3.alignment = PP_ALIGN.CENTER

p4 = tf2.add_paragraph()
p4.text = "KU Leuven"
p4.font.size = Pt(20)
p4.font.color.rgb = RGBColor(150, 190, 230)
p4.alignment = PP_ALIGN.CENTER


# ============================================================
# SLIDE 2: AGENDA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_title(slide, "Agenda")

items = [
    "1.  Problem Setup & Cost Structure",
    "2.  Data Exploration & Key Patterns",
    "3.  Demand Forecasting Model",
    "4.  Inventory Simulation Engine",
    "5.  Ordering Policy Design",
    "6.  Backtest Results & Performance",
    "7.  Live Game Demo",
]
add_body_text(slide, 1.5, 1.8, 10, 5, items, font_size=22)


# ============================================================
# SLIDE 3: PROBLEM SETUP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_title(slide, "Problem Setup", "Lost-sales model with perishable inventory")

# Cost table
data = [
    ["Parameter", "Value", "Impact"],
    ["Holding cost", "1 / unit / period", "Cheap to hold"],
    ["Shortage cost", "19 / unit", "Very expensive to stockout"],
    ["Expiry cost", "9 / unit", "Moderate waste penalty"],
    ["Lead time", "2 periods", "Orders arrive in 2 days"],
    ["Shelf life", "2 periods", "FIFO: oldest sold first"],
]
add_table(slide, 0.8, 1.6, 6.5, 3.5, data, col_widths=[2.0, 2.5, 2.0])

# Key insight box
add_accent_box(slide, 7.8, 1.8, 4.8, 1.2,
               "Critical Ratio = 19/20 = 0.95\nTarget the 95th percentile of demand",
               bg_color=RGBColor(255, 245, 230), font_color=ACCENT_ORANGE, font_size=18)

# Tension box
add_accent_box(slide, 7.8, 3.4, 4.8, 1.6,
               "The Core Tension:\nShortage (19) says over-order\nExpiry (9) says don't waste\nShelf life (2) limits buffer",
               bg_color=LIGHT_BG, font_color=DARK_BLUE, font_size=16)

# Starting position
add_body_text(slide, 0.8, 5.4, 11, 1.5, [
    ("Starting position [5, 4, 3]: ", "5 units arriving tomorrow, 4 fresh on-hand, 3 expiring today = 7 on-hand + 5 pipeline"),
    ("Game: ", "26 periods (Jul 2 - Aug 2, 2021), excluding Sundays & holidays"),
], font_size=16)


# ============================================================
# SLIDE 4: SIMULATION SEQUENCE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_title(slide, "Simulation Sequence", "What happens each period")

# 4 boxes for the 4 steps
steps = [
    ("1. Start of Period", "Pipeline arrives\nObserve inventory\nDecide order (arrives t+2)"),
    ("2. During Period", "Demand occurs\nFIFO: sell oldest first\nUnmet demand = shortage"),
    ("3. End of Period", "Age-1 units expire (cost 9)\nHolding cost on remaining\nAge-0 becomes age-1"),
    ("4. Next Period", "New shipments arrive\nRepeat"),
]

for i, (title, body) in enumerate(steps):
    left = 0.6 + i * 3.15
    # Title box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.8), Inches(2.9), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Body box
    shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.6), Inches(2.9), Inches(2.0))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = LIGHT_BG
    shape2.line.color.rgb = MED_BLUE
    shape2.line.width = Pt(1)
    tf2 = shape2.text_frame
    tf2.word_wrap = True
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    for j, line in enumerate(body.split("\n")):
        p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)

# Arrow between boxes
for i in range(3):
    left = 3.4 + i * 3.15
    txBox = slide.shapes.add_textbox(Inches(left), Inches(2.9), Inches(0.4), Inches(0.6))
    tf = txBox.text_frame
    tf.paragraphs[0].text = ">"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.color.rgb = ACCENT_ORANGE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# State representation
add_body_text(slide, 0.8, 5.2, 11.5, 1.5, [
    ("State tracked: ", "on_hand = [age_0, age_1]  |  pipeline = [arriving_next, arriving_in_2]"),
    ("Key constraint: ", "With lead time 2, orders placed today serve demand 2 days from now"),
], font_size=16)


# ============================================================
# SLIDE 5: DATA EXPLORATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_title(slide, "Data Exploration", "SKU 2921141  |  1,165 training days  |  Sep 2017 - Jul 2021")

# Left: key stats
add_body_text(slide, 0.8, 1.6, 5.5, 4.5, [
    ("Sales statistics:", ""),
    ("  Mean: 5.0", "  |  Median: 4  |  Std: 4.3  |  Max: 24"),
    "",
    ("Promotion effect:", ""),
    ("  With promo: ", "5.77 avg  (+33% uplift)"),
    ("  Without promo: ", "4.33 avg"),
    "",
    ("Day-of-week effect:", ""),
    ("  Saturday: ", "~2x weekday demand (strongest signal)"),
    ("  Friday: ", "~1.4x weekday demand"),
    "",
    ("Price sensitivity:", ""),
    ("  Low price (<2.5): ", "mean ~6.3 (independent of promo)"),
], font_size=16)

# Right: test period structure
add_accent_box(slide, 7.5, 1.6, 5.2, 0.6, "Test Period Structure (26 days)", font_size=18)

test_data = [
    ["Phase", "Dates", "Promo", "Exp. Demand"],
    ["Phase 1", "Jul 2-13 (12d)", "Yes (40%)", "~6-8 / day"],
    ["Phase 2", "Jul 14-Aug 2 (14d)", "No", "~3-5 / day"],
]
add_table(slide, 7.5, 2.5, 5.2, 1.5, test_data, col_widths=[1.0, 1.8, 1.2, 1.2])

add_accent_box(slide, 7.5, 4.4, 5.2, 1.8,
               "Key risk: Promo > Non-promo\ntransition on Jul 14\n\nDemand drops suddenly\nExcess inventory may expire",
               bg_color=RGBColor(255, 240, 240), font_color=RED, font_size=15)


# ============================================================
# SLIDE 6: DEMAND FORECASTING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_title(slide, "Demand Forecasting", "LightGBM Quantile Regression")

# Left: approach
add_body_text(slide, 0.8, 1.6, 5.5, 2.5, [
    ("Why quantile regression?", ""),
    "  - Need the full demand distribution, not just a point forecast",
    "  - Directly answers: 'What is the 95th percentile of demand?'",
    "  - Handles feature interactions (DOW x promo x price)",
    "  - Trains in < 1 second, robust with 1,165 rows",
    "",
    ("5 models trained:", ""),
    "  - Poisson mean model (reference)",
    "  - Q50, Q75, Q90, Q95 (quantile regression)",
], font_size=15)

# Right: features
add_accent_box(slide, 7.0, 1.6, 5.8, 0.5, "8 Features", font_size=16)
feat_data = [
    ["Feature", "Type", "Why"],
    ["dow (0-6)", "Engineered", "Weekly seasonality"],
    ["is_saturday", "Binary", "Highest demand day"],
    ["is_friday", "Binary", "2nd highest demand"],
    ["PROMO_01", "Binary", "Promotion active"],
    ["PROMO_DEPTH", "0-60%", "Discount magnitude"],
    ["PRC_2_norm", "Continuous", "Price sensitivity"],
    ["HOLIDAY_f1", "Binary", "Tomorrow is holiday"],
    ["HOLIDAY_l1", "Binary", "Yesterday was holiday"],
]
add_table(slide, 7.0, 2.3, 5.8, 3.6, feat_data, col_widths=[1.8, 1.5, 2.5])

# Validation results
add_accent_box(slide, 0.8, 4.5, 5.5, 0.5, "Validation (last 200 days holdout)", font_size=16)
val_data = [
    ["Quantile", "Pinball Loss", "Coverage", "Target"],
    ["Q50", "1.552", "64.0%", "50%"],
    ["Q75", "1.340", "79.5%", "75%"],
    ["Q90", "0.803", "90.0%", "90%"],
    ["Q95", "0.504", "93.0%", "95%"],
]
add_table(slide, 0.8, 5.2, 5.5, 1.6, val_data, col_widths=[1.2, 1.5, 1.3, 1.0])


# ============================================================
# SLIDE 7: ORDERING POLICY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_title(slide, "Ordering Policy", "Perishable Newsvendor with Forward Inventory Projection")

# 3-step process
step_titles = ["Step 1: Forecast", "Step 2: Project Forward", "Step 3: Compute Order"]
step_bodies = [
    "Generate quantile forecasts\nfor periods t, t+1, t+2, t+3\nusing known features\n(DOW, promo, price, holidays)",
    "Simulate FIFO + expiry\nthrough t and t+1 to estimate\nhow much inventory survives\nto period t+2",
    "order = max(0,\n  Q95(t+2) - existing_at_t+2)\n\nCapped by max useful demand\nFloored by median demand",
]

for i in range(3):
    left = 0.6 + i * 4.2
    # Step number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 1.3), Inches(1.5), Inches(0.6), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_ORANGE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title
    txBox = slide.shapes.add_textbox(Inches(left), Inches(2.2), Inches(3.8), Inches(0.5))
    tf2 = txBox.text_frame
    tf2.paragraphs[0].text = step_titles[i]
    tf2.paragraphs[0].font.size = Pt(20)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = DARK_BLUE
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Body
    shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.8), Inches(3.8), Inches(2.4))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = LIGHT_BG
    shape2.line.color.rgb = MED_BLUE
    shape2.line.width = Pt(1)
    tf3 = shape2.text_frame
    tf3.word_wrap = True
    tf3.vertical_anchor = MSO_ANCHOR.MIDDLE
    for j, line in enumerate(step_bodies[i].split("\n")):
        p = tf3.paragraphs[0] if j == 0 else tf3.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)

# End-game tapering box
add_accent_box(slide, 0.8, 5.6, 5.5, 1.2,
               "End-Game Tapering\n>3 periods left: Q95 | 3: Q90 | 2: Q75 | <=1: Q50",
               bg_color=RGBColor(255, 245, 230), font_color=ACCENT_ORANGE, font_size=14)

# Why it works box
add_accent_box(slide, 6.8, 5.6, 5.8, 1.2,
               "Why this works:\nReactive to actual demand via state updates\nFeature-aware (adapts to promo/DOW/price)\nBalances shortage vs expiry via quantile + caps",
               bg_color=LIGHT_BG, font_color=DARK_BLUE, font_size=13)


# ============================================================
# SLIDE 8: FORWARD PROJECTION DETAIL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_title(slide, "Forward Inventory Projection", "How we estimate what survives to t+2")

# Timeline diagram using boxes
periods = [
    ("Period t (now)", "Pipeline[0] arrives\nas age-0\n\nDemand consumes\nage-1 first (FIFO)\n\nAge-1 remainder\nexpires\n\nAge-0 survivors\ncarry forward"),
    ("Period t+1", "Pipeline[1] arrives\nas age-0\n\nCarry-forward from t\nbecomes age-1\n\nDemand consumes\nage-1 first (FIFO)\n\nAge-0 survivors\ncarry forward"),
    ("Period t+2", "NEW ORDER arrives\nas age-0\n\nCarry-forward from t+1\nbecomes age-1\n(= existing_at_t2)\n\nOrder fills the gap\nto Q95 target"),
]

for i, (title, body) in enumerate(periods):
    left = 0.6 + i * 4.2
    color = ACCENT_ORANGE if i == 2 else DARK_BLUE

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.6), Inches(3.8), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.3), Inches(3.8), Inches(3.8))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = LIGHT_BG
    shape2.line.color.rgb = color
    shape2.line.width = Pt(1.5)
    tf2 = shape2.text_frame
    tf2.word_wrap = True
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    for j, line in enumerate(body.split("\n")):
        p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(2)

# Arrows
for i in range(2):
    left = 4.3 + i * 4.2
    txBox = slide.shapes.add_textbox(Inches(left), Inches(3.5), Inches(0.5), Inches(0.6))
    tf = txBox.text_frame
    tf.paragraphs[0].text = ">>"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = ACCENT_ORANGE
    tf.paragraphs[0].font.bold = True

# Formula at bottom
add_accent_box(slide, 2.5, 6.3, 8.3, 0.6,
               "order = max(0,  Q_target(t+2)  -  existing_at_t2 )",
               bg_color=RGBColor(255, 245, 230), font_color=ACCENT_ORANGE, font_size=20)


# ============================================================
# SLIDE 9: BACKTEST RESULTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_title(slide, "Backtest Results", "Policy performance under different demand scenarios")

# Results table
data = [
    ["Scenario", "Total Cost", "Holding", "Shortage", "Expiry"],
    ["Median demand", "375", "119", "76", "180"],
    ["Q75 demand", "639", "12", "627", "0"],
    ["Q90 demand (stress)", "1,522", "2", "1,520", "0"],
]
add_table(slide, 1.5, 1.6, 10.3, 2.0, data, col_widths=[2.5, 1.8, 1.8, 1.8, 1.8])

# Interpretation boxes
add_accent_box(slide, 0.8, 4.0, 3.8, 2.5,
               "Median scenario\n(most likely)\n\nTotal: 375\nSmall unavoidable P2 shortage\nMost expiry at promo transition\nand end-game",
               bg_color=RGBColor(235, 250, 240), font_color=GREEN, font_size=14)

add_accent_box(slide, 4.9, 4.0, 3.8, 2.5,
               "Key observations\n\nP2 shortage is structural\n(pipeline gap [5,0])\n\nZero expiry in high-demand\nscenarios = no over-ordering\n\nPolicy adapts in live play\nvia state updates",
               bg_color=LIGHT_BG, font_color=DARK_BLUE, font_size=14)

add_accent_box(slide, 9.0, 4.0, 3.8, 2.5,
               "Live game advantage\n\nActual demand revealed\neach period\n\nHigh demand depletes stock\n-> higher next order\n\nLow demand builds buffer\n-> lower next order",
               bg_color=RGBColor(255, 245, 230), font_color=ACCENT_ORANGE, font_size=14)


# ============================================================
# SLIDE 10: ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_title(slide, "System Architecture", "Three modular Python files")

# Three component boxes
components = [
    ("demand_model.py", "Demand Forecasting",
     "LightGBM quantile regression\n5 models (mean, Q50-Q95)\n8 engineered features\nTrains in < 1 second"),
    ("inventory_engine.py", "Inventory Simulation",
     "FIFO + expiry mechanics\nAge cohort tracking\nFull cost computation\nComplete period history"),
    ("play_game.py", "Live Game Interface",
     "Interactive CLI\nRecommended order each period\nForward projection policy\nEnd-game tapering"),
]

for i, (filename, title, body) in enumerate(components):
    left = 0.6 + i * 4.2

    # Filename
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.6), Inches(3.8), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = filename
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title
    shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.3), Inches(3.8), Inches(0.5))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = MED_BLUE
    shape2.line.fill.background()
    tf2 = shape2.text_frame
    tf2.paragraphs[0].text = title
    tf2.paragraphs[0].font.size = Pt(16)
    tf2.paragraphs[0].font.color.rgb = WHITE
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Body
    shape3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.9), Inches(3.8), Inches(2.2))
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = LIGHT_BG
    shape3.line.color.rgb = MED_BLUE
    shape3.line.width = Pt(1)
    tf3 = shape3.text_frame
    tf3.word_wrap = True
    tf3.vertical_anchor = MSO_ANCHOR.MIDDLE
    for j, line in enumerate(body.split("\n")):
        p = tf3.paragraphs[0] if j == 0 else tf3.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)

# Flow arrows
for i in range(2):
    left = 4.3 + i * 4.2
    txBox = slide.shapes.add_textbox(Inches(left), Inches(3.2), Inches(0.5), Inches(0.6))
    tf = txBox.text_frame
    tf.paragraphs[0].text = ">>"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = ACCENT_ORANGE
    tf.paragraphs[0].font.bold = True

# Bottom workflow
add_accent_box(slide, 1.5, 5.5, 10.3, 0.8,
               "Workflow:  Train models once (2s)  >>  Each period: forecast demand, project inventory, recommend order  >>  User confirms or overrides",
               bg_color=RGBColor(255, 245, 230), font_color=ACCENT_ORANGE, font_size=15)


# ============================================================
# SLIDE 11: KEY DESIGN DECISIONS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_title(slide, "Key Design Decisions")

decisions = [
    ("LightGBM Quantile > Parametric Distributions: ",
     "Data is overdispersed (var/mean = 3.66) with complex interactions. "
     "Quantile regression directly outputs percentiles without distributional assumptions."),
    ("Myopic Newsvendor > Reinforcement Learning: ",
     "26 periods + 2-period look-ahead = small problem. Analytical policy is near-optimal, "
     "trains instantly, and runs fast during live play."),
    ("Single-SKU Training > Pooled Model: ",
     "1,165 rows is sufficient for 8 features. Other subgroup-109 SKUs have different demand levels "
     "(2.9 and 6.9 vs 5.0) - pooling would blur the signal."),
    ("Cost Asymmetry Drives Everything: ",
     "Shortage (19) >> Holding (1). One stockout = 19 days of holding. "
     "Policy deliberately over-orders, capped only by shelf-life expiry risk."),
    ("End-Game Tapering: ",
     "Progressively reduce target quantile (Q95 > Q90 > Q75 > Q50) as game ends "
     "to prevent wasteful final-period expiry."),
]

for i, (bold_part, normal_part) in enumerate(decisions):
    top = 1.6 + i * 1.05
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.7), Inches(0.95))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = bold_part
    run1.font.bold = True
    run1.font.size = Pt(16)
    run1.font.color.rgb = DARK_BLUE
    run2 = p.add_run()
    run2.text = normal_part
    run2.font.size = Pt(15)
    run2.font.color.rgb = DARK_GRAY


# ============================================================
# SLIDE 12: THANK YOU / Q&A
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(48)
p.font.color.rgb = WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Questions & Discussion"
p2.font.size = Pt(28)
p2.font.color.rgb = LIGHT_BLUE
p2.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(8), Inches(2))
tf2 = txBox2.text_frame
summary_items = [
    "Feature-driven quantile forecasting with LightGBM",
    "Perishable newsvendor policy with forward FIFO/expiry projection",
    "Adaptive end-game tapering to minimize waste",
    "Interactive CLI for fast live-game decisions",
]
for i, item in enumerate(summary_items):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.text = f"   {item}"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(180, 210, 240)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(8)


# ============================================================
# SAVE
# ============================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
